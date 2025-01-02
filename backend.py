from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt, Inches
import cohere
import os

def summaryGenerate(prompt,complexity):
  co = cohere.Client(
      #api_key =
  )
  complexityInt = int(complexity)
  complexity_append = ""
  if complexityInt == 1:
     complexity_append = "at a Flesch-Kincaid reading comprehension level of 60 (meant for upper middle school student reading comprehension) "
  elif complexityInt == 2:
     complexity_append = "at a Flesch-Kincaid reading comprehension level of 30 (meant for upper high school student reading comprehension) "
  else:
     complexity_append = "at a Flesch-Kincaid reading comprehension level of 0 (meant for PhD student reading comprehension) "
  chat = co.chat(
      message= "Create an ONLY 5 (no more or less) bullet summary of this article "+complexity_append+"using dashes (-) to seperate each point: "+prompt,
      model="command"
  )
  lines = chat.text.split("\n")
  bullet_points = []
  start_bullet_points = False
  for line in lines:
      if line.strip().startswith("-"):
          bullet_points.append(line.strip())
          start_bullet_points = True
      elif start_bullet_points and not line.strip():
          break
  array = [item[2:] for item in bullet_points]
  return array

def citationConvert(citation,desiredFormat):
  co = cohere.Client(
      #api_key =
  )
  desiredFormatInt = int(desiredFormat)
  format_append = ""
  if desiredFormatInt == 1:
     format_append = "IEEE Citation Format"
  elif desiredFormatInt == 2:
     format_append = "APA Citation Format"
  else:
     format_append = "MLA Citation Format"
  chat = co.chat(
      message= "Convert the following citation into "+format_append+" and do not give any other information or reccomendations: "+citation,
      model="command"
  )
  lines = chat.text.split("\n")
  return lines[2].replace('*','"')

def prescreate(article_title,aiArray,cite,col):
  if col == "B":
     backgr = RGBColor(0,0,0)
     fontCol = RGBColor(255,255,0)
     pptName = "Black_Yellow_Slides.pptx"
  else:
     backgr = RGBColor(255,255,255)
     fontCol = RGBColor(0,0,0)
     pptName = "White_Black_Slides.pptx"
  prs = Presentation()
  title_slide_layout = prs.slide_layouts[0]
  slide = prs.slides.add_slide(title_slide_layout)
  background = slide.background
  fill = background.fill
  fill.solid()
  fill.fore_color.rgb = backgr
  title = slide.shapes.title
  subtitle = slide.placeholders[1]
  title.text = article_title
  subtitle.text = "A Brief Overview"
  image_path = "sickkids.png"
  slide.shapes.add_picture(image_path,Inches(7.3), Inches(6.6))
  for shape in slide.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = fontCol
  
  slide_layout = prs.slide_layouts[1]
  slide = prs.slides.add_slide(slide_layout)
  title = slide.shapes.title
  background = slide.background
  fill = background.fill
  fill.solid()
  fill.fore_color.rgb = backgr
  title.text = "Summary Notes"
  bullet_slide = slide.placeholders[1]
  bullet_slide.text_frame.margin_top = Inches(-0.3)
  bullet_points = aiArray
  for point in bullet_points:
      paragraph = bullet_slide.text_frame.add_paragraph()
      paragraph.text = point
      paragraph.font.size = Inches(0.25)
  slide.shapes.add_picture(image_path,Inches(7.3), Inches(6.6))
  for shape in slide.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = fontCol
  
  slide_layout = prs.slide_layouts[1]
  slide = prs.slides.add_slide(slide_layout)
  title = slide.shapes.title
  background = slide.background
  fill = background.fill
  fill.solid()
  fill.fore_color.rgb = backgr
  title.text = "References"
  bullet_slide = slide.placeholders[1]
  bullet_slide.text_frame.margin_top = Inches(-0.3)
  bullet_points = [cite]
  for point in bullet_points:
      paragraph = bullet_slide.text_frame.add_paragraph()
      paragraph.text = point
      paragraph.font.size = Inches(0.25)
  slide.shapes.add_picture(image_path,Inches(7.3), Inches(6.6))
  for shape in slide.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = fontCol
  prs.save(pptName)

print("Welcome to Team 30's backend prototype! Enter the title of the literature review topic below:")
article_title = input()
print("\nEnter the complexity level of the powerpoint summary (1,2,3):\n1. Upper-Middle School level\n2. Upper-High School level\n3. University level")
complexity = input()
print("\nEnter article as a single line (visit https://lingojam.com/TexttoOneLine to convert for free):")
article = input()
print("\nEnter article citation in any format as a single line (visit https://lingojam.com/TexttoOneLine to convert for free):")
citation = input()
print("\nEnter desired format for article citation (1,2,3):\n1. IEEE\n2. APA\n3. MLA")
desiredFormat = input()
text = summaryGenerate(article,complexity)
cite = citationConvert(citation,desiredFormat)
prescreate(article_title,text,cite,"W")
prescreate(article_title,text,cite,"B")
print("\nPPT Generation Complete! You can find the PPTS in the dist folder")
